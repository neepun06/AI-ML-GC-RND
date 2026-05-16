"""Reusable slide components: header, footer, container, metric tile, bullets.

All visual primitives the deck renderer composes from. No business logic here —
just shapes and text styled per `theme.py`.
"""
from __future__ import annotations

from typing import Iterable

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from kelp_teaser.render import theme
from kelp_teaser.schemas.slide import Bullet, MetricTile


def _style_run(run, *, font=None, size=None, bold=False, color=None):
    run.font.name = font or theme.BODY_FONT
    if size is not None:
        run.font.size = size
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_header(slide, *, codename: str, subtitle: str = "") -> None:
    """Draw the Dark Indigo top band with codename + subtitle + Kelp logo placeholder."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                theme.SLIDE_W, theme.HEADER_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.PALETTE["primary"]
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = codename
    _style_run(p.runs[0], font=theme.HEADING_FONT, size=theme.HEADING_SIZE,
               bold=True, color=theme.PALETTE["white"])

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(9), Inches(0.3))
        p2 = sub_box.text_frame.paragraphs[0]
        p2.text = subtitle
        _style_run(p2.runs[0], font=theme.BODY_FONT, size=theme.SUBHEADING_SIZE,
                   color=theme.PALETTE["white"])

    # Kelp logo placeholder (text, top-right). Real logo asset wired later if desired.
    logo_box = slide.shapes.add_textbox(Inches(11.6), Inches(0.35), Inches(1.4), Inches(0.6))
    p3 = logo_box.text_frame.paragraphs[0]
    p3.text = theme.LOGO_PLACEHOLDER_TEXT
    p3.alignment = PP_ALIGN.RIGHT
    _style_run(p3.runs[0], font=theme.HEADING_FONT, size=theme.HEADING_SIZE,
               bold=True, color=theme.PALETTE["white"])


def add_footer(slide) -> None:
    """Brand-mandated centered footer at 9pt."""
    box = slide.shapes.add_textbox(Inches(0), Inches(7.15), Inches(13.3), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = theme.FOOTER_TEXT
    p.alignment = PP_ALIGN.CENTER
    _style_run(p.runs[0], size=theme.FOOTER_FONT_SIZE, color=theme.PALETTE["text_muted"])


def draw_container(slide, x: Emu, y: Emu, w: Emu, h: Emu, *, title: str = ""):
    """White rounded card with a subtle shadow. Returns the card shape."""
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Pt(3), y + Pt(3), w, h)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = theme.PALETTE["border_light"]
    shadow.line.fill.background()

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = theme.PALETTE["white"]
    box.line.color.rgb = theme.PALETTE["border_light"]

    if title:
        tb = slide.shapes.add_textbox(x + Pt(10), y + Pt(8), w - Pt(20), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = title.upper()
        _style_run(p.runs[0], font=theme.HEADING_FONT, size=Pt(11),
                   bold=True, color=theme.PALETTE["primary"])

    return box


def draw_metric_tile(slide, x: Emu, y: Emu, w: Emu, h: Emu, tile: MetricTile) -> None:
    """A big-number metric card. Value on top, label below, optional subtext."""
    draw_container(slide, x, y, w, h)

    val_box = slide.shapes.add_textbox(x, y + Pt(5), w, Inches(0.7))
    p = val_box.text_frame.paragraphs[0]
    p.text = tile.value
    p.alignment = PP_ALIGN.CENTER
    _style_run(p.runs[0], font=theme.HEADING_FONT, size=Pt(28),
               bold=True, color=theme.PALETTE["accent"])

    lbl_box = slide.shapes.add_textbox(x, y + Pt(50), w, Inches(0.3))
    p = lbl_box.text_frame.paragraphs[0]
    p.text = tile.label.upper()
    p.alignment = PP_ALIGN.CENTER
    _style_run(p.runs[0], font=theme.HEADING_FONT, size=Pt(9),
               bold=True, color=theme.PALETTE["primary"])

    if tile.subtext:
        sub_box = slide.shapes.add_textbox(x, y + Pt(70), w, Inches(0.4))
        p = sub_box.text_frame.paragraphs[0]
        p.text = tile.subtext
        p.alignment = PP_ALIGN.CENTER
        _style_run(p.runs[0], size=Pt(9), color=theme.PALETTE["text_muted"])


def draw_bullet_list(slide, x: Emu, y: Emu, w: Emu, h: Emu,
                     bullets: Iterable[Bullet]) -> None:
    """Stack of body-text bullets inside the given rectangle."""
    tb = slide.shapes.add_textbox(x + Pt(8), y + Pt(8), w - Pt(16), h - Pt(16))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"• {b.text}"
        p.space_after = Pt(6)
        _style_run(p.runs[0], size=Pt(11), color=theme.PALETTE["text_dark"])
