"""Native python-pptx chart renderer for the 6 supported chart kinds."""
from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Emu

from kelp_teaser.schemas.plan import ChartKind
from kelp_teaser.schemas.slide import ChartSpec

_KIND_TO_XL: dict[ChartKind, int] = {
    ChartKind.revenue_growth_bar: XL_CHART_TYPE.COLUMN_CLUSTERED,
    ChartKind.revenue_growth_line: XL_CHART_TYPE.LINE,
    ChartKind.segment_mix_donut: XL_CHART_TYPE.DOUGHNUT,
    ChartKind.margin_trend_line: XL_CHART_TYPE.LINE,
    ChartKind.geo_split_stacked_bar: XL_CHART_TYPE.COLUMN_STACKED,
    ChartKind.channel_mix_donut: XL_CHART_TYPE.DOUGHNUT,
}


def render_chart(slide, x: Emu, y: Emu, w: Emu, h: Emu, spec: ChartSpec) -> None:
    """Render a native PowerPoint chart onto `slide` per `spec`."""
    if spec is None or not isinstance(spec, ChartSpec):
        raise ValueError("render_chart requires a ChartSpec instance")

    xl_kind = _KIND_TO_XL.get(spec.chart_kind)
    if xl_kind is None:
        raise ValueError(f"unsupported chart kind: {spec.chart_kind!r}")

    data = CategoryChartData()
    data.categories = spec.categories
    for series in spec.series:
        data.add_series(series.name, series.values)

    # python-pptx serialises chart graphicFrame x/y/cx/cy with str(value), so a
    # float EMU like 4389120.0 lands in the XML as "4389120.0". PowerPoint then
    # refuses to parse the graphicFrame and strips the chart on repair. Force
    # ints here so callers can do EMU arithmetic with `/` without breaking us.
    chart = slide.shapes.add_chart(
        xl_kind, int(x), int(y), int(w), int(h), data,
    ).chart
    chart.has_title = bool(spec.title)
    if spec.title:
        chart.chart_title.text_frame.text = spec.title

    # Value data labels on bar/line charts (donuts get legend instead).
    if spec.chart_kind not in (ChartKind.segment_mix_donut,
                               ChartKind.channel_mix_donut):
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.number_format = "General"
        plot.data_labels.number_format_is_linked = False

    if spec.chart_kind in (ChartKind.segment_mix_donut, ChartKind.channel_mix_donut):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
