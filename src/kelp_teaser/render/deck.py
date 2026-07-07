"""Render a list of ComposedSlide objects into a .pptx file.

The renderer is dumb on purpose: it only knows how to draw what each
ComposedSection says. All content decisions happen upstream in the agents.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from kelp_teaser.render import theme
from kelp_teaser.render.charts import render_chart
from kelp_teaser.render.slide_components import (
    add_footer,
    add_header,
    add_picture_cover,
    draw_bullet_list,
    draw_container,
    draw_metric_tile,
)
from kelp_teaser.schemas.plan import ComponentKind
from kelp_teaser.schemas.slide import ComposedSection, ComposedSlide


# Height weight per section kind. Charts/images/quadrants need real vertical
# room; metric/kpi strips are short. Bullets are medium.
_ROW_WEIGHTS: dict[ComponentKind, int] = {
    ComponentKind.chart: 5,
    ComponentKind.hero_image: 5,
    ComponentKind.quadrant: 4,
    ComponentKind.bullet_list: 2,
    ComponentKind.product_grid: 2,
    ComponentKind.metric_tile: 1,
    ComponentKind.kpi_strip: 1,
}


def render_deck(*, slides: list[ComposedSlide], codename: str, out_path: Path) -> Path:
    """Render the deck. Returns the saved path."""
    indices = [s.index for s in slides]
    assert sorted(indices) == list(range(len(slides))), (
        f"Slide indices must be contiguous 0..{len(slides) - 1}, "
        f"got {sorted(indices)}"
    )

    prs = Presentation()
    prs.slide_width = theme.SLIDE_W
    prs.slide_height = theme.SLIDE_H

    for composed in sorted(slides, key=lambda s: s.index):
        _render_slide(prs, composed, codename)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def _render_slide(prs, composed: ComposedSlide, codename: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = theme.PALETTE["bg_slide"]

    add_header(slide, codename=codename, subtitle=composed.title)
    add_footer(slide)

    content_h = theme.SLIDE_H - theme.HEADER_HEIGHT - Inches(0.6)

    # Build render units: a [chart, bullet_list] adjacent pair is ONE unit
    # (rendered side-by-side) weighted by the chart; everything else is its
    # own unit weighted by its kind.
    secs = composed.sections
    render_weights: list[int] = []
    j = 0
    while j < len(secs):
        if (secs[j].kind == ComponentKind.chart
                and j + 1 < len(secs)
                and secs[j + 1].kind == ComponentKind.bullet_list):
            render_weights.append(_ROW_WEIGHTS[ComponentKind.chart])
            j += 2
        else:
            render_weights.append(_ROW_WEIGHTS.get(secs[j].kind, 2))
            j += 1
    total_weight = sum(render_weights) or 1

    y_cursor = theme.HEADER_HEIGHT + Inches(0.2)
    unit_idx = 0
    i = 0
    while i < len(secs):
        sec = secs[i]
        nxt = secs[i + 1] if i + 1 < len(secs) else None
        row_h = content_h * render_weights[unit_idx] / total_weight
        # Composite: chart + adjacent bullet_list -> side-by-side band.
        if (sec.kind == ComponentKind.chart
                and nxt is not None
                and nxt.kind == ComponentKind.bullet_list):
            half_w = (theme.CONTENT_W - theme.GUTTER) / 2
            _render_section(slide, sec, x=theme.MARGIN, y=y_cursor,
                            w=half_w, h=row_h - Inches(0.1))
            _render_section(slide, nxt,
                            x=theme.MARGIN + half_w + theme.GUTTER, y=y_cursor,
                            w=half_w, h=row_h - Inches(0.1))
            i += 2
        else:
            _render_section(slide, sec, x=theme.MARGIN, y=y_cursor,
                            w=theme.CONTENT_W, h=row_h - Inches(0.1))
            i += 1
        y_cursor += row_h
        unit_idx += 1


def _render_section(slide, section: ComposedSection, *, x, y, w, h) -> None:
    if section.kind == ComponentKind.bullet_list:
        draw_container(slide, x, y, w, h, title=section.heading)
        draw_bullet_list(slide, x, y + Inches(0.4), w, h - Inches(0.4), section.bullets)

    elif section.kind == ComponentKind.metric_tile:
        # Equal horizontal split across N tiles.
        n = max(len(section.metrics), 1)
        tile_w = (w - (n - 1) * theme.GUTTER) / n
        cx = x
        for tile in section.metrics:
            draw_metric_tile(slide, cx, y, tile_w, h, tile)
            cx += tile_w + theme.GUTTER

    elif section.kind == ComponentKind.chart and section.chart is not None:
        # The chart's own title carries the heading, so we skip the container
        # title and keep the vertical inset small to preserve plot height.
        draw_container(slide, x, y, w, h)
        render_chart(slide,
                     x + Inches(0.2), y + Inches(0.15),
                     w - Inches(0.4), h - Inches(0.3),
                     section.chart)

    elif section.kind == ComponentKind.hero_image and section.image is not None:
        # If the image file actually exists on disk, render it. Otherwise
        # (Composer hallucinated a path, or ImageCurator was skipped because
        # PEXELS_API_KEY is not set) fail soft to a labelled placeholder so
        # the deck still renders.
        image_path = Path(section.image.path)
        if image_path.is_file():
            add_picture_cover(slide, str(image_path), x, y, w, h)
        else:
            placeholder_title = (section.heading or section.image.alt_text
                                 or "Image unavailable")
            draw_container(slide, x, y, w, h, title=placeholder_title)

    elif section.kind == ComponentKind.product_grid:
        draw_container(slide, x, y, w, h, title=section.heading or "Portfolio")
        draw_bullet_list(slide, x, y + Inches(0.4), w, h - Inches(0.4), section.bullets)

    elif section.kind == ComponentKind.kpi_strip:
        # Render as horizontal metric tiles (alias of metric_tile layout).
        n = max(len(section.metrics), 1)
        tile_w = (w - (n - 1) * theme.GUTTER) / n
        cx = x
        for tile in section.metrics:
            draw_metric_tile(slide, cx, y, tile_w, h, tile)
            cx += tile_w + theme.GUTTER

    elif section.kind == ComponentKind.quadrant:
        # 2×2 bullet quadrants. Caller must pass 4 bullets to fill all quadrants.
        draw_container(slide, x, y, w, h, title=section.heading)
        half_w = (w - theme.GUTTER) / 2
        half_h = (h - theme.GUTTER - Inches(0.4)) / 2
        cy0 = y + Inches(0.4)
        positions = [
            (x, cy0),
            (x + half_w + theme.GUTTER, cy0),
            (x, cy0 + half_h + theme.GUTTER),
            (x + half_w + theme.GUTTER, cy0 + half_h + theme.GUTTER),
        ]
        for (px, py), bullet in zip(positions, section.bullets[:4]):
            draw_bullet_list(slide, px, py, half_w, half_h, [bullet])
