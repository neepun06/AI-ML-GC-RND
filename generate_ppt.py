import json
import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 1. SETUP
TEMPLATE_PATH = "Master_Template.pptx"
OUTPUT_FOLDER = "Final_Submissions"

if not os.path.exists(OUTPUT_FOLDER):
    os.makedirs(OUTPUT_FOLDER)

def replace_text(shape, placeholder, replacement):
    """Recursively replaces text in shapes (and groups)."""
    if not shape.has_text_frame:
        return
    
    # Simple replacement
    if placeholder in shape.text:
        # Preserve formatting by replacing inside the run if possible
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, str(replacement))

def update_chart(slide, chart_data_json):
    """Finds the first chart on the slide and updates its data."""
    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            
            # Create new chart data
            new_chart_data = CategoryChartData()
            new_chart_data.categories = chart_data_json['years']
            
            # Add Revenue Series
            new_chart_data.add_series('Revenue', chart_data_json['revenue'])
            
            # Push data to chart
            chart.replace_data(new_chart_data)
            print("   📊 Chart data updated!")
            return

def generate_presentation(company_name):
    json_path = f"{company_name}_ANALYSIS.json"
    
    # 1. Load Data
    if not os.path.exists(json_path):
        print(f"❌ JSON not found for {company_name}. Run analyze.py first.")
        return

    with open(json_path, 'r') as f:
        data = json.load(f)
        
    print(f"🎨 Generating PPT for {company_name}...")
    try:
        prs = Presentation(TEMPLATE_PATH)
    except Exception as e:
        print(f"❌ Template error: {e}. Make sure 'Master_Template.pptx' is closed!")
        return

    # --- SLIDE 1: OVERVIEW ---
    slide1 = prs.slides[0]
    # Map JSON keys to your Template Placeholders
    replacements_1 = {
        "{{TEASER_TITLE}}": data['slide_1'].get('teaser_title', 'Investment Teaser'),
        "{{PROJECT_NAME}}": data['slide_1'].get('project_name', 'Project Confidential'),
        "{{SECTOR}}": data['slide_1'].get('sector', 'General'),
    }
    
    for shape in slide1.shapes:
        for key, val in replacements_1.items():
            replace_text(shape, key, val)
            
    # Add Bullet Points (Special Handling)
    # Find the shape meant for "Business Summary" (You might need to identify it by ID or position in real usage)
    # For hackathon speed, we assume there's a shape containing "{{SUMMARY_BULLETS}}"
    for shape in slide1.shapes:
        if "{{SUMMARY_BULLETS}}" in shape.text:
            shape.text_frame.clear() # Clear the placeholder
            for bullet in data['slide_1']['business_summary']:
                p = shape.text_frame.add_paragraph()
                p.text = bullet
                p.level = 0

    # --- SLIDE 2: FINANCIALS ---
    slide2 = prs.slides[1]
    replacements_2 = {
        "{{CAGR}}": str(data['slide_2'].get('revenue_cagr', 'N/A')),
        "{{EBITDA}}": str(data['slide_2'].get('ebitda_margin', 'N/A')),
    }
    for shape in slide2.shapes:
        for key, val in replacements_2.items():
            replace_text(shape, key, val)
            
    # Update Chart
    if 'graph_data' in data['slide_2']:
        update_chart(slide2, data['slide_2']['graph_data'])

    # --- SLIDE 3: THESIS ---
    slide3 = prs.slides[2]
    # Hooks
    hooks = data['slide_3'].get('investment_thesis', [])
    hooks_text = "\n".join([f"• {h}" for h in hooks])
    
    for shape in slide3.shapes:
        replace_text(shape, "{{INVESTMENT_THESIS}}", hooks_text)

    # --- SAVE ---
    output_path = os.path.join(OUTPUT_FOLDER, f"{company_name}_Teaser.pptx")
    prs.save(output_path)
    print(f"✅ Saved to {output_path}")

# --- TEST ---
if __name__ == "__main__":
    generate_presentation("Test Company Name")