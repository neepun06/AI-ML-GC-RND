from pptx import Presentation
from pptx.util import Pt

from kelp_teaser.render import theme
from kelp_teaser.render.deck import render_deck
from kelp_teaser.schemas.slide import Bullet, ComposedSection, ComposedSlide
from kelp_teaser.schemas.plan import ComponentKind


def _slide(i):
    return ComposedSlide(index=i, title=f"S{i}", sections=[
        ComposedSection(kind=ComponentKind.bullet_list, bullets=[
            Bullet(text="x", source_id="doc:x.md")]),
    ])


def test_every_slide_has_footer_and_logo(tmp_path):
    out = render_deck(slides=[_slide(0), _slide(1), _slide(2)],
                      codename="Project X", out_path=tmp_path / "d.pptx")
    prs = Presentation(str(out))
    for slide in prs.slides:
        texts = [sh.text_frame.text for sh in slide.shapes
                 if sh.has_text_frame]
        assert any(theme.FOOTER_TEXT in t for t in texts), "footer missing"
        assert any(theme.LOGO_PLACEHOLDER_TEXT == t.strip() for t in texts), \
            "logo placeholder missing"


def test_footer_is_9pt():
    assert theme.FOOTER_FONT_SIZE == Pt(9)


def test_fonts_are_arial():
    assert theme.HEADING_FONT == "Arial"
    assert theme.BODY_FONT == "Arial"
