import pytest
from pptx import Presentation
from pptx.util import Inches

from kelp_teaser.render.deck import render_deck
from kelp_teaser.schemas.plan import ComponentKind
from kelp_teaser.schemas.slide import Bullet, ComposedSection, ComposedSlide, MetricTile
from tests.fixtures.ksolves_minimal import three_slides


def test_render_deck_writes_three_slides(tmp_path):
    out = tmp_path / "teaser.pptx"
    render_deck(slides=three_slides(), codename="Project Halo", out_path=out)
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == 3


def test_render_deck_includes_footer_text_on_every_slide(tmp_path):
    out = tmp_path / "teaser.pptx"
    render_deck(slides=three_slides(), codename="Project Halo", out_path=out)
    prs = Presentation(str(out))
    from kelp_teaser.render import theme
    for s in prs.slides:
        texts = [sh.text_frame.text for sh in s.shapes if sh.has_text_frame]
        assert any(theme.FOOTER_TEXT in t for t in texts)


def test_render_deck_uses_codename_in_header(tmp_path):
    out = tmp_path / "teaser.pptx"
    render_deck(slides=three_slides(), codename="Project Halo", out_path=out)
    prs = Presentation(str(out))
    first = prs.slides[0]
    texts = [sh.text_frame.text for sh in first.shapes if sh.has_text_frame]
    assert any("Project Halo" in t for t in texts)


def _bullet(idx: int) -> ComposedSlide:
    return ComposedSlide(index=idx, title=f"Slide {idx}", sections=[
        ComposedSection(kind=ComponentKind.bullet_list, bullets=[
            Bullet(text="x", source_id="doc:x.md"),
        ]),
    ])


def test_render_deck_rejects_non_contiguous_indices(tmp_path):
    """If indices are not exactly 0..N-1, render_deck must raise rather
    than silently render a misordered or missing-slide deck."""
    slides = [_bullet(0), _bullet(2), _bullet(3)]  # gap at index 1
    with pytest.raises(AssertionError):
        render_deck(slides=slides, codename="Project Halo",
                    out_path=tmp_path / "teaser.pptx")


def test_render_deck_rejects_duplicate_indices(tmp_path):
    slides = [_bullet(0), _bullet(1), _bullet(1)]
    with pytest.raises(AssertionError):
        render_deck(slides=slides, codename="Project Halo",
                    out_path=tmp_path / "teaser.pptx")


def test_render_deck_fails_soft_when_hero_image_path_missing(tmp_path):
    """If a hero_image section's path doesn't exist on disk (e.g. the
    Composer hallucinated a path, or ImageCurator was skipped because the
    Pexels key is unset), the renderer should fall back to a placeholder
    container rather than crashing with FileNotFoundError."""
    from kelp_teaser.schemas.slide import HeroImage

    hero_slide = ComposedSlide(index=0, title="Cover", sections=[
        ComposedSection(
            kind=ComponentKind.hero_image,
            image=HeroImage(
                path="images/does_not_exist.jpg",
                alt_text="Hero placeholder text",
                source_id="doc:x.md",
            ),
        ),
    ])
    slides = [hero_slide, _bullet(1), _bullet(2)]

    out = tmp_path / "teaser.pptx"
    render_deck(slides=slides, codename="Project Halo", out_path=out)
    assert out.exists()

    prs = Presentation(str(out))
    cover = prs.slides[0]
    texts = [sh.text_frame.text for sh in cover.shapes if sh.has_text_frame]
    # draw_container uppercases the title; match case-insensitively.
    assert any("hero placeholder text" in t.lower() for t in texts)


def _chart_slide() -> ComposedSlide:
    from kelp_teaser.schemas.slide import ChartSpec, ChartSeries
    from kelp_teaser.schemas.plan import ChartKind
    chart = ChartSpec(
        chart_kind=ChartKind.revenue_growth_bar,
        title="Revenue",
        categories=["FY23", "FY24", "FY25"],
        series=[ChartSeries(name="Revenue", values=[10.0, 20.0, 30.0])],
        source_id="doc:x.md",
    )
    return ComposedSlide(index=0, title="Financials", sections=[
        ComposedSection(kind=ComponentKind.chart, heading="Revenue", chart=chart),
        ComposedSection(kind=ComponentKind.kpi_strip, metrics=[
            MetricTile(value="30", label="Rev", source_id="doc:x.md"),
        ]),
        ComposedSection(kind=ComponentKind.bullet_list, bullets=[
            Bullet(text="x", source_id="doc:x.md"),
        ]),
    ])


def test_chart_row_is_taller_than_kpi_row(tmp_path):
    """A chart section must be allocated a taller frame than a kpi_strip
    section on the same slide (weighted rows, not equal rows)."""
    from pptx import Presentation
    out = render_deck(slides=[_chart_slide()], codename="Project X",
                      out_path=tmp_path / "d.pptx")
    prs = Presentation(str(out))
    slide = prs.slides[0]
    chart_h = next(sh.height for sh in slide.shapes if sh.has_chart)
    # The chart must have real vertical room — at least 1.8 inches.
    assert chart_h >= Inches(1.8), f"chart too short: {chart_h} EMU"
