from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from kelp_teaser.render import theme
from kelp_teaser.render.slide_components import (
    add_header,
    add_footer,
    add_picture_cover,
    draw_container,
    draw_metric_tile,
    draw_bullet_list,
)
from kelp_teaser.schemas.slide import Bullet, MetricTile


def _blank_slide():
    prs = Presentation()
    prs.slide_width = theme.SLIDE_W
    prs.slide_height = theme.SLIDE_H
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_add_header_adds_shapes():
    prs, slide = _blank_slide()
    before = len(slide.shapes)
    add_header(slide, codename="Project Halo", subtitle="Manufacturing | Teaser")
    assert len(slide.shapes) > before


def test_add_footer_text_matches_brand():
    prs, slide = _blank_slide()
    add_footer(slide)
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    assert any(theme.FOOTER_TEXT in t for t in texts)


def test_draw_container_returns_shape():
    prs, slide = _blank_slide()
    shape = draw_container(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(3), title="Profile")
    assert shape is not None


def test_draw_metric_tile_renders_value_and_label():
    prs, slide = _blank_slide()
    tile = MetricTile(value="₹450 Cr", label="Revenue FY24", source_id="doc:r.pdf#p1")
    draw_metric_tile(slide, Inches(0.5), Inches(1.5), Inches(2.0), Inches(1.8), tile)
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    assert any("₹450 Cr" in t for t in texts)
    assert any("REVENUE FY24" in t.upper() for t in texts)


def test_draw_bullet_list_renders_each_bullet():
    prs, slide = _blank_slide()
    bullets = [
        Bullet(text="5 facilities in western India", source_id="doc:r.pdf#p1"),
        Bullet(text="600+ active customers", source_id="web:tavily:https://x.com"),
    ]
    draw_bullet_list(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(4), bullets)
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    joined = "\n".join(texts)
    assert "5 facilities" in joined
    assert "600+ active" in joined


def _make_wide_image(path: Path) -> None:
    # 400x100 — a 4:1 wide image.
    Image.new("RGB", (400, 100), (120, 120, 200)).save(path)


def test_add_picture_cover_preserves_aspect_ratio(tmp_path):
    """A 4:1 image placed in a 2:1 box must keep its 4:1 aspect ratio
    (scaled to cover), not be squashed to the box ratio."""
    img = tmp_path / "wide.png"
    _make_wide_image(img)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Box is 6in x 3in (2:1).
    pic = add_picture_cover(slide, str(img),
                            Inches(1), Inches(1), Inches(6), Inches(3))
    ratio = pic.width / pic.height
    assert abs(ratio - 4.0) < 0.05, f"aspect not preserved: {ratio}"
