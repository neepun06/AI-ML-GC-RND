from pptx import Presentation
from pptx.util import Inches

from kelp_teaser.render import theme
from kelp_teaser.render.charts import render_chart
from kelp_teaser.schemas.slide import ChartSeries, ChartSpec


def _blank_slide():
    prs = Presentation()
    prs.slide_width = theme.SLIDE_W
    prs.slide_height = theme.SLIDE_H
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_render_revenue_bar_chart():
    prs, slide = _blank_slide()
    spec = ChartSpec(
        chart_kind="revenue_growth_bar",
        title="Revenue",
        categories=["FY22", "FY23", "FY24"],
        series=[ChartSeries(name="Revenue (₹ Cr)", values=[300, 380, 450])],
        source_id="doc:r.pdf#p12",
    )
    render_chart(slide, Inches(1), Inches(2), Inches(8), Inches(4), spec)
    chart_shapes = [s for s in slide.shapes if s.has_chart]
    assert len(chart_shapes) == 1


def test_render_segment_mix_donut():
    prs, slide = _blank_slide()
    spec = ChartSpec(
        chart_kind="segment_mix_donut",
        title="Segment mix",
        categories=["Lecithin", "Phospholipids", "Other"],
        series=[ChartSeries(name="Share", values=[55, 30, 15])],
        source_id="doc:r.pdf#p15",
    )
    render_chart(slide, Inches(1), Inches(2), Inches(5), Inches(4), spec)
    chart_shapes = [s for s in slide.shapes if s.has_chart]
    assert len(chart_shapes) == 1


def test_render_unknown_chart_kind_raises():
    import pytest
    prs, slide = _blank_slide()

    # Construct invalid spec via dict (bypass enum) — should fail upstream, but
    # render_chart itself should also reject if somehow reached.
    with pytest.raises(ValueError):
        render_chart(slide, Inches(1), Inches(2), Inches(5), Inches(4),
                     spec=None)  # type: ignore[arg-type]


def test_render_chart_coerces_float_emu_to_int():
    """python-pptx serialises chart graphicFrame coordinates with str(value),
    so a float EMU like 4389120.0 would land in slide XML as "4389120.0" and
    cause PowerPoint to strip the chart on repair. render_chart must coerce
    to int so callers using `/` (which produces floats) stay safe."""
    import re
    import zipfile

    prs, slide = _blank_slide()
    spec = ChartSpec(
        chart_kind="revenue_growth_line",
        title="Revenue",
        categories=["2020", "2021", "2022"],
        series=[ChartSeries(name="Revenue", values=[100, 200, 300])],
        source_id="doc:x.md",
    )
    # Float EMU values (what deck.py's row_h math produces).
    render_chart(slide, x=914400.0, y=4389120.0,
                 w=10911535.0, h=1920240.0, spec=spec)

    import io
    buf = io.BytesIO()
    prs.save(buf)
    with zipfile.ZipFile(io.BytesIO(buf.getvalue())) as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml").decode("utf-8")
    # No graphicFrame attribute should look like "4389120.0".
    bad = re.findall(r'(x|y|cx|cy)="[0-9]+\.[0-9]+"', slide_xml)
    assert bad == [], (
        f"Found non-integer EMU coords in slide XML: {bad}. "
        "PowerPoint will refuse to load this and strip the chart on repair."
    )


def test_bar_chart_has_data_labels(tmp_path):
    """Column/line charts must show value data labels so numbers are
    readable without reading the axis."""
    from pptx import Presentation
    from pptx.util import Inches
    from kelp_teaser.render.charts import render_chart
    from kelp_teaser.schemas.slide import ChartSpec, ChartSeries
    from kelp_teaser.schemas.plan import ChartKind

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec = ChartSpec(
        chart_kind=ChartKind.revenue_growth_bar, title="Rev",
        categories=["FY23", "FY24"],
        series=[ChartSeries(name="Revenue", values=[10.0, 20.0])],
        source_id="doc:x.md",
    )
    render_chart(slide, Inches(1), Inches(1), Inches(5), Inches(3), spec)
    chart = next(sh.chart for sh in slide.shapes if sh.has_chart)
    plot = chart.plots[0]
    assert plot.has_data_labels is True
