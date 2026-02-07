import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from utils import download_image 

# ==========================================
# 1. DESIGN SYSTEM (Kelp Branding)
# ==========================================
THEME = {
    "primary": RGBColor(40, 0, 80),      # Dark Indigo
    "accent": RGBColor(255, 100, 80),    # Pink/Orange Gradient
    "bg_slide": RGBColor(245, 247, 250), # Light SaaS Grey
    "white": RGBColor(255, 255, 255),
    "text_dark": RGBColor(45, 45, 55),
    "text_muted": RGBColor(100, 100, 110),
    "border_light": RGBColor(220, 220, 230)
}

""" dimensions and typography are defined in one place for consistency and easy adjustments """

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.5)
GUTTER = Inches(0.3)

CONTENT_W = SLIDE_W - (MARGIN * 2)

COL2_W = (CONTENT_W - GUTTER) / 2
COL3_W = (CONTENT_W - GUTTER * 2) / 3

class TextStyler:
    """Centralized Typography Control"""
    @staticmethod
    def title(run):
        run.font.name = "Arial"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = THEME["white"]

    @staticmethod
    def subtitle(run):
        run.font.name = "Arial"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(200, 200, 255)

    @staticmethod
    def card_header(run):
        run.font.name = "Arial"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = THEME["primary"]
        
    @staticmethod
    def body(run, size=10, bold=False, color=None):
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color if color else THEME["text_dark"]

    @staticmethod
    def metric_big(run):
        run.font.name = "Arial"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = THEME["accent"]

# ==========================================
# 2. DRAWING PRIMITIVES
# ==========================================

def add_header(slide, company_name, sector, tagline=None):
    """Draws the Dark Indigo Hero Header"""
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.3))
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME["primary"]
    bg.line.fill.background()
    
    # Text
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = company_name
    TextStyler.title(p.runs[0] if p.runs else p.add_run())
    
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(9), Inches(0.3))
    p2 = tb2.text_frame.paragraphs[0]
    final_tag = tagline if tagline else f"{sector} | Investment Teaser"
    p2.text = final_tag
    TextStyler.subtitle(p2.runs[0] if p2.runs else p2.add_run())


    
    if os.path.exists("logo.png"):
        slide.shapes.add_picture("logo.png", Inches(11.6), Inches(0.35), height=Inches(0.6))

def add_footer(slide):
    box = slide.shapes.add_textbox(Inches(0), Inches(7.15), Inches(13.3), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = "Strictly Private & Confidential – Prepared by Kelp M&A Team"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.color.rgb = THEME["text_muted"]

def draw_container(slide, x, y, w, h, title=None):
    """Draws a white card with shadow and optional title"""
    # Shadow
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x+Pt(3), y+Pt(3), w, h)
    s.fill.solid() # <--- FIXED: Added .solid()
    s.fill.fore_color.rgb = RGBColor(210,210,210)
    s.line.fill.background()
    
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid() # <--- FIXED: Added .solid()
    box.fill.fore_color.rgb = THEME["white"]
    box.line.color.rgb = THEME["border_light"]
    
    if title:
        tb = slide.shapes.add_textbox(x+Pt(10), y+Pt(8), w-Pt(20), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = title.upper()
        TextStyler.card_header(p.runs[0] if p.runs else p.add_run())
    
    return box

def draw_stat_tile(slide, x, y, w, h, value, label, subtext=""):
    """Draws a High-Density Metric Tile (Gamma Style)"""
    box = draw_container(slide, x, y, w, h)
    
    # Value (Big Number)
    tb_val = slide.shapes.add_textbox(x, y+Pt(5), w, Inches(0.7))
    p = tb_val.text_frame.paragraphs[0]
    p.text = str(value)
    p.alignment = PP_ALIGN.CENTER
    TextStyler.metric_big(p.runs[0] if p.runs else p.add_run())

    # Label (Medium)
    tb_lbl = slide.shapes.add_textbox(x, y+Pt(50), w, Inches(0.3))
    p = tb_lbl.text_frame.paragraphs[0]
    p.text = label.upper()
    p.alignment = PP_ALIGN.CENTER
    TextStyler.card_header(p.runs[0] if p.runs else p.add_run())
    p.runs[0].font.size = Pt(9)

    # Context (Small Grey)
    if subtext:
        tb_sub = slide.shapes.add_textbox(x, y+Pt(70), w, Inches(0.4))
        p = tb_sub.text_frame.paragraphs[0]
        p.text = subtext
        p.alignment = PP_ALIGN.CENTER
        TextStyler.body(p.runs[0] if p.runs else p.add_run(), size=9, color=THEME["text_muted"])

# ==========================================
# 3. SLIDE BUILDERS (The Logic)
# ==========================================

def build_slide_1(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = THEME["bg_slide"]
    
    s1 = data.get('slide_1', {})
    add_header(slide, s1.get('project_name',''), s1.get('sector',''), s1.get('header_tagline'))
    add_footer(slide)

    # 1. Business Overview (Top Left)
    draw_container(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(3.5), "Business Profile")
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(5.6), Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True
    
    # Narrative
    overview = s1.get('business_overview_card', {})
    if overview.get('main_narrative'):
        p = tf.add_paragraph()
        p.text = overview['main_narrative']
        TextStyler.body(p.runs[0] if p.runs else p.add_run(), size=11, bold=True)
        p.space_after = Pt(10)
    
    # Bullets
    for bullet in overview.get('bullets', []):
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.space_after = Pt(6)
        TextStyler.body(p.runs[0] if p.runs else p.add_run(), size=10.5)

    # 2. Infrastructure Highlights (Row of 3 Tiles below Overview)
    y_stats = Inches(5.2)
    stats = s1.get('infrastructure_highlights', [])
    # Fallback if AI didn't generate stats
    if not stats: stats = [{"value": "N/A", "label": "Metrics", "subtext": ""}]
        
    for i, stat in enumerate(stats[:3]):
        x = Inches(0.5) + (i * Inches(2.1))
        draw_stat_tile(slide, x, y_stats, Inches(2.0), Inches(1.8), 
                       stat.get('value',''), stat.get('label',''), stat.get('subtext',''))

    # 3. Key Offerings Grid (Right Side)
    # This replaces the empty space with a structured Product Grid
    draw_container(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5.5), "Key Portfolio & Capabilities")
    
    products = s1.get('key_products_grid', [])
    y_prod = Inches(2.0)
    
    for prod in products[:5]: # Max 5 items
        # Product Box
        tb = slide.shapes.add_textbox(Inches(7.0), y_prod, Inches(5.6), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = prod.get('name', 'Product').upper()
        TextStyler.card_header(p.runs[0] if p.runs else p.add_run())
        
        # Spec/Detail
        tb2 = slide.shapes.add_textbox(Inches(7.0), y_prod+Pt(15), Inches(5.6), Inches(0.5))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        p.text = prod.get('desc', '')
        TextStyler.body(p.runs[0] if p.runs else p.add_run(), size=10, color=RGBColor(80,80,90))
        
        y_prod += Inches(0.9)


def build_slide_2(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = THEME["bg_slide"]
    
    s2 = data.get('slide_2', {})
    add_header(slide, "Financial Performance", "Operational Scale", "Growth Trajectory")
    add_footer(slide)

    # 1. Top Row: 3 Major Financial Cards
    fin_card = s2.get('financial_card', {})
    metrics = [
        ("Revenue CAGR", fin_card.get('cagr_value', 'N/A'), "FY20-24"),
        ("EBITDA Margin", fin_card.get('ebitda_value', 'N/A'), "Sustainable"),
        ("FY25E Revenue", fin_card.get('revenue_fy25', 'N/A'), "Projected")
    ]
    
    for i, (lbl, val, sub) in enumerate(metrics):
        draw_stat_tile(slide, Inches(0.5 + i*4.2), Inches(1.5), Inches(4), Inches(1.5), val, lbl, sub)

    # 2. Bottom Left: Native Chart
    draw_container(slide, Inches(0.5), Inches(3.3), Inches(8), Inches(3.7), "Revenue Growth (INR Cr)")
    
    chart_json = s2.get('graph_data', {})
    chart_data = CategoryChartData()
    chart_data.categories = chart_json.get('years', ['23', '24', '25'])
    vals = [float(x) for x in chart_json.get('revenue', [10, 20, 30])]
    chart_data.add_series('Revenue', vals)

    x, y, cx, cy = Inches(0.7), Inches(3.8), Inches(7.5), Inches(3.0)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.has_title = False

    # 3. Bottom Right: Operational Grid
    draw_container(slide, Inches(8.8), Inches(3.3), Inches(4), Inches(3.7), "Operational KPIs")
    
    op_grid = s2.get('operational_kpi_grid', [])
    y_kpi = Inches(3.8)
    
    for kpi in op_grid[:4]:
        # Metric
        tb = slide.shapes.add_textbox(Inches(9.0), y_kpi, Inches(3.6), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = f"{kpi.get('metric', '')} : {kpi.get('label', '')}"
        TextStyler.card_header(p.runs[0] if p.runs else p.add_run())
        
        # Context
        tb2 = slide.shapes.add_textbox(Inches(9.0), y_kpi+Pt(15), Inches(3.6), Inches(0.4))
        p = tb2.text_frame.paragraphs[0]
        p.text = kpi.get('context', '')
        TextStyler.body(p.runs[0] if p.runs else p.add_run(), size=10, color=THEME["text_muted"])
        
        y_kpi += Inches(0.8)


def build_slide_3(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = THEME["bg_slide"]
    
    s3 = data.get('slide_3', {})
    add_header(slide, "Investment Highlights", "Strategic Rationale", "Why Invest?")
    add_footer(slide)

    hooks = s3.get('thesis_points', [])
    y = Inches(1.6)
    
    for i, hook in enumerate(hooks):
        # Card
        draw_container(slide, Inches(1), y, Inches(11.3), Inches(1.4))
        
        # Number Circle
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y+Inches(0.45), Inches(0.5), Inches(0.5))
        oval.fill.solid()
        oval.fill.fore_color.rgb = THEME["accent"]
        oval.line.fill.background()
        
        # Title
        tb = slide.shapes.add_textbox(Inches(1.3), y+Pt(10), Inches(10), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = hook.get('title', 'Highlight').upper()
        TextStyler.card_header(p.runs[0] if p.runs else p.add_run())
        p.runs[0].font.size = Pt(12)
        p.runs[0].font.color.rgb = THEME["primary"]
        
        # Detail
        tb2 = slide.shapes.add_textbox(Inches(1.3), y+Pt(35), Inches(10.8), Inches(0.8))
        tf = tb2.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = hook.get('detail', '')
        TextStyler.body(p.runs[0] if p.runs else p.add_run(), size=11)
        
        y += Inches(1.7)

def build_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = THEME["bg_slide"]

    add_header(slide, "Important Notice & Disclaimer", "")
    add_footer(slide)

    draw_container(slide, MARGIN, Inches(1.6),
                   CONTENT_W, Inches(5.5))

    disclaimer = (
        "Strictly Private & Confidential. Prepared exclusively for the intended recipient. "
        "This presentation is informational and not an offer to sell securities. "
        "Information believed reliable but not guaranteed. Forward-looking statements "
        "and projections may vary materially from actual results."
    )

    tb = slide.shapes.add_textbox(
        MARGIN + Inches(0.3),
        Inches(2.3),
        CONTENT_W - Inches(0.6),
        Inches(4.8)
    )

    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = disclaimer
    TextStyler.body(p.runs[0] if p.runs else p.add_run(), size=19)




# ==========================================
# 4. MAIN RUNNER
# ==========================================
def generate_deck(company_name):
    json_path = f"{company_name}_ANALYSIS.json"
    print(f"🚀 PPT Engine: Assembling High-Density Deck for {company_name}...")
    
    if not os.path.exists(json_path):
        print(f"❌ Critical Error: Data file '{json_path}' not found. Run analyze.py first.")
        return

    with open(json_path, 'r') as f:
        data = json.load(f)

    prs = Presentation()
    # Force Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    try:
        build_slide_1(prs, data)
        build_slide_2(prs, data)
        build_slide_3(prs, data)
        build_slide_4(prs)
        
        out_path = f"Final_Submissions/{company_name}_Teaser_Atomic.pptx"
        if not os.path.exists("Final_Submissions"):
            os.makedirs("Final_Submissions")
            
        prs.save(out_path)
        print(f"✅ Success! High-Density Deck Generated: {out_path}")
        
    except Exception as e:
        print(f"❌ Build Failed: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    generate_deck("Test Company Name")